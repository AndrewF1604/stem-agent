"""Generuje prezentację PowerPoint całego projektu (LLM Error Analysis).

Uruchom z głównego folderu projektu:
    python build_pptx.py
Wynik: wyniki_i_wnioski/Prezentacja_LLM_Error_Analysis.pptx
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ACCENT = RGBColor(0x1F, 0x38, 0x64)
DARK = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
IMG = "wyniki_i_wnioski/"
OUT = "wyniki_i_wnioski/Prezentacja_LLM_Error_Analysis.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def title_bar(s, text):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.16))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.32), SW - Inches(1.0), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = ACCENT


def _rich(p, text, size, color=DARK):
    for part in re.split(r"(\*\*.*?\*\*)", text):
        if not part:
            continue
        run = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]; run.font.bold = True
        else:
            run.text = part
        run.font.size = Pt(size); run.font.color.rgb = color


def bullets(s, items, left, top, width, height, size=18):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        lead = p.add_run(); lead.text = "•  "
        lead.font.size = Pt(size); lead.font.color.rgb = ACCENT; lead.font.bold = True
        _rich(p, item, size)


def table(s, headers, rows, left, top, width, height, fs=14):
    gt = s.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for j, h in enumerate(headers):
        c = gt.cell(0, j); c.text = h
        run = c.text_frame.paragraphs[0].runs[0]
        run.font.bold = True; run.font.size = Pt(fs); run.font.color.rgb = WHITE
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = str(val)
            for para in c.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(fs); run.font.color.rgb = DARK


def note(s, text, top=Inches(6.7)):
    tb = s.shapes.add_textbox(Inches(0.5), top, SW - Inches(1.0), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    _rich(p, text, 13, GRAY)


# ---------- Slajd tytułowy ----------
s = new_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), SW, Inches(2.0))
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.8), SW - Inches(1.6), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Gdzie poddaje się model 8B?"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); r2 = p2.add_run()
r2.text = "Analiza błędów Llama 3.1 8B w multi-hop QA (HotpotQA)"
r2.font.size = Pt(22); r2.font.color.rgb = WHITE
sub = s.shapes.add_textbox(Inches(0.8), Inches(4.9), SW - Inches(1.6), Inches(1.2))
ps = sub.text_frame.paragraphs[0]
_rich(ps, "Llama 3.1 8B  ·  7 konfiguracji  ·  sędzia Gemma 2 9B", 18, GRAY)
ps2 = sub.text_frame.add_paragraph()
_rich(ps2, "Kurs NLP  ·  [imiona autorów]", 16, GRAY)

# ---------- 1. Metodologia ----------
s = new_slide(); title_bar(s, "Problem i metodologia")
bullets(s, [
    "**Pytania multi-hop** — odpowiedź wymaga *dwóch skoków* (np. reżyser filmu → jego miejsce urodzenia)",
    "**Przepływ:** 100 zamrożonych pytań → 7 konfiguracji → 3400 odpowiedzi → 3 metryki + sędzia AI",
    "**Zbiór:** HotpotQA (distractor), 100 pytań zbalansowanych (25× bridge/comparison × easy/hard), zamrożony raz",
    "**Model:** Llama 3.1 8B lokalnie (Ollama)",
    "**7 konfiguracji:** temperatura (0 / 0.5 / 1.0) × prompt (standard / CoT / expert) × agent z Wikipedią",
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.0), size=20)

# ---------- 2. Kryzys metryk ----------
s = new_slide(); title_bar(s, "Kryzys klasycznych metryk")
bullets(s, [
    "**Exact Match = 0%** we wszystkich 7 konfiguracjach",
    "**F1 = 3–17%**",
    "**Sędzia AI (Gemma 2 9B) = 69–77%**",
    "Model „rozlewa\" odpowiedź: gold *Elizabeth Taylor*, model *„Dame Elizabeth Rosemond Taylor was a British-American actress…\"*",
    "**Wniosek:** miary n-gramowe są bezużyteczne do oceny wnioskowania LLM",
], Inches(0.5), Inches(1.7), Inches(6.6), Inches(5.0), size=17)
s.shapes.add_picture(IMG + "performance_by_config.png", Inches(7.2), Inches(2.0), width=Inches(5.7))

# ---------- 3. Typologia ----------
s = new_slide(); title_bar(s, "Typologia błędów")
bullets(s, [
    "**Multi-hop reasoning — 80.9%**",
    "Entity confusion — 8.2%",
    "Format error — 5.6%",
    "Halucynacja — 4.5%",
    "Negacja — 0.8%",
    "**Wniosek:** bariera 8B to drugi skok rozumowania, nie halucynacje",
], Inches(0.5), Inches(1.7), Inches(6.4), Inches(5.0), size=18)
s.shapes.add_picture(IMG + "error_distribution.png", Inches(7.1), Inches(1.9), width=Inches(5.8))

# ---------- 4. Sys vs stoch ----------
s = new_slide(); title_bar(s, "Błędy systematyczne vs stochastyczne")
bullets(s, [
    "**temp 0** (C1, C4, C6) → błędy **systematyczne** (0 stochastycznych — determinizm)",
    "**temp 0.5–1.0** (C2, C3, C5) → błędy **stochastyczne** rosną z temperaturą (C3: 62/100)",
    "**Recepta:** self-consistency naprawia stochastyczne, **nie** systematyczne",
    "**Teza:** ustawienia zmieniają nie ILE, a JAK (i czy naprawialnie) model się myli",
], Inches(0.5), Inches(1.7), Inches(6.4), Inches(5.0), size=17)
s.shapes.add_picture(IMG + "systematic_stochastic.png", Inches(7.1), Inches(1.9), width=Inches(5.8))

# ---------- 5. Paradoks RAG ----------
s = new_slide(); title_bar(s, "Paradoks RAG: agent z Wikipedią wypada GORZEJ")
bullets(s, [
    "C7 (agent + Wikipedia): **~28%** sędziego",
    "vs **69–77%** przy kontekście w prompcie",
    "Błędy C7 są **systematyczne** (72/100) — powtarzalnie złe trajektorie",
    "**Dla modeli 8B „podaj kontekst\" bije „pozwól szukać\"**",
    "Model 8B gubi się w pętli ReAct: kiepskie zapytania, błądzenie",
], Inches(0.5), Inches(1.7), Inches(6.6), Inches(5.0), size=17)
s.shapes.add_picture(IMG + "performance_by_config.png", Inches(7.2), Inches(2.0), width=Inches(5.7))

# ---------- 6. Wpływ promptu ----------
s = new_slide(); title_bar(s, "Wpływ promptu")
table(s, ["Konfiguracja", "Judge"], [
    ["C6 — expert (temp 0)", "0.77  (najlepszy)"],
    ["C4 — CoT (temp 0)", "0.75"],
    ["C1 — standard (temp 0)", "0.73"],
    ["C3 — standard (temp 1.0)", "0.69  (najgorszy kontekstowy)"],
], Inches(0.7), Inches(1.9), Inches(7.0), Inches(2.6), fs=16)
bullets(s, [
    "CoT / expert dają **+2–4 pkt** — nie rewolucję",
    "CoT **zmienia profil błędów**: F1 spada do ~3% (model jeszcze bardziej „leje wodę\")",
    "Wysoka temperatura lekko szkodzi",
], Inches(0.7), Inches(4.8), Inches(11.5), Inches(2.2), size=18)

# ---------- 7. Walidacja kappa ----------
s = new_slide(); title_bar(s, "Walidacja ludzka: Cohen's kappa i Gwet's AC1")
bullets(s, [
    "**Surowa zgodność człowiek ↔ Gemma: 74%**",
    "**Cohen's kappa = 0.256**   ·   **Gwet's AC1 = 0.715** (znaczna zgodność)",
    "**Paradoks kappy:** rozkład skośny (76–84% multi-hop) zaniża κ; AC1 (odporny, Gwet 2008) pokazuje zgodność jako znaczną",
    "**Realny finding:** człowiek wyłapał **15 błędów formatu** (model dobrze, ale rozwlekle), które Gemma wzięła za multi-hop",
    "→ Sędzia AI **myli formę z poprawnością** (jak EM) — zautomatyzowany pipeline wymaga nadzoru człowieka",
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.0), size=18)

# ---------- 8. Wnioski ----------
s = new_slide(); title_bar(s, "Wnioski")
bullets(s, [
    "**1.** Miary n-gramowe (EM/F1) są martwe dla oceny LLM — liczy się sędzia semantyczny",
    "**2.** Bariera modeli 8B to **drugi skok rozumowania** (80.9% błędów), nie halucynacje",
    "**3.** RAG bez silnego modelu szkodzi; self-consistency naprawia tylko błędy stochastyczne",
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), size=19)
bullets(s, [
    "**Ograniczenia:** 100 pytań (25/komórkę) · jeden anotator · sędzia 9B · możliwa kontaminacja HotpotQA",
    "**Future work:** silniejszy sędzia (GPT-4o/Gemini) · self-consistency produkcyjnie · lepszy scaffolding agenta",
], Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.3), size=16)

# ---------- 9. Dziękujemy ----------
s = new_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.0), SW, Inches(1.5))
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.25), SW - Inches(1.6), Inches(1.1))
p = tb.text_frame.paragraphs[0]; r = p.add_run()
r.text = "Dziękujemy — pytania?"
r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = WHITE
note(s, "Slajdy backup na kolejnych stronach: pełna tabela metryk, taksonomia, przykład multi-hop, liczby sys/stoch, sanity-check EM.", Inches(5.0))

# ---------- BACKUP 1: tabela 7x3 ----------
s = new_slide(); title_bar(s, "BACKUP — Pełna tabela metryk (7×3)")
table(s, ["Konfiguracja", "EM", "F1", "Judge"], [
    ["C1 — std, temp 0", "0%", "16.7%", "73.0%"],
    ["C2 — std, temp 0.5 (×10)", "0%", "17.1%", "73.5%"],
    ["C3 — std, temp 1.0 (×10)", "0%", "15.1%", "69.1%"],
    ["C4 — CoT, temp 0", "0%", "3.5%", "75.0%"],
    ["C5 — CoT, temp 0.5 (×10)", "0%", "3.3%", "74.0%"],
    ["C6 — expert, temp 0", "0%", "5.7%", "77.0%"],
    ["C7 — agent + Wikipedia", "0%", "4.0%", "28.0%"],
], Inches(1.2), Inches(1.8), Inches(10.9), Inches(4.4), fs=15)
note(s, "F1 spada przy CoT/expert (C4–C6), bo dłuższe „rozumowanie\" rozcieńcza pokrycie tokenów.")

# ---------- BACKUP 2: taksonomia ----------
s = new_slide(); title_bar(s, "BACKUP — Taksonomia 5 kategorii błędów")
bullets(s, [
    "**1. Halucynacja** — informacja zmyślona, nieobecna w tekście",
    "**2. Multi-hop reasoning** — pierwszy skok OK, drugi krok zawodzi (połączenie faktów)",
    "**3. Format** — zna fakt, ale zła forma (akapit „lania wody\" zamiast krótkiej frazy)",
    "**4. Negacja** — model źle odczytał negację w tekście źródłowym",
    "**5. Entity confusion** — wybrał powiązaną, ale złą encję",
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.6), size=19)
note(s, "Rozkład: multi-hop 80.9% · entity 8.2% · format 5.6% · halucynacja 4.5% · negacja 0.8%")

# ---------- BACKUP 3: przyklad multi-hop ----------
s = new_slide(); title_bar(s, "BACKUP — Przykład: czysty multi-hop failure")
bullets(s, [
    "**Pytanie:** *Charles Libby graduated from a high school located at what address?*",
    "**Złota odpowiedź:** 284 Cumberland Avenue",
    "**Model (C1, temp 0):** *„There is no information… However, he graduated from Portland High School in Maine.\"*",
    "**Fakt ①:** *Libby graduated from Portland High School* — ✅ model to znalazł",
    "**Fakt ②:** *Portland High School is located at 284 Cumberland Avenue* — ❌ model nie połączył",
    "→ Pierwszy skok wykonany, drugi nie. Model nawet **przyznaje** „brak informacji\", choć była w tekście.",
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.0), size=17)

# ---------- BACKUP 4: sys/stoch liczby ----------
s = new_slide(); title_bar(s, "BACKUP — Systematyczne vs stochastyczne (n=100)")
table(s, ["Konfiguracja", "Bez błędu", "Stochastyczne", "Systematyczne"], [
    ["C1 — temp 0", "73", "0", "27"],
    ["C2 — temp 0.5", "49", "51", "0"],
    ["C3 — temp 1.0", "38", "62", "0"],
    ["C4 — CoT temp 0", "75", "0", "25"],
    ["C5 — CoT temp 0.5", "28", "72", "0"],
    ["C6 — expert temp 0", "77", "0", "23"],
    ["C7 — agent", "28", "0", "72"],
], Inches(1.2), Inches(1.8), Inches(10.9), Inches(4.4), fs=15)
note(s, "temp 0 → wyłącznie systematyczne (determinizm). temp>0 → stochastyczne rosną z temperaturą.")

# ---------- BACKUP 5: sanity EM ----------
s = new_slide(); title_bar(s, "BACKUP — Sanity-check: dlaczego EM = 0%")
bullets(s, [
    "**EM = 0 dla wszystkich 3400 odpowiedzi** — realny efekt, nie bug",
    "**Mediana odpowiedzi modelu: 39 słów** vs **mediana golda: 2 słowa** (~20× za długo)",
    "**Złota odpowiedź dosłownie zawarta w 68.2%** odpowiedzi → fakt zwykle JEST, metryka go gubi",
    "**„Zawiera gold\" 68.2% ≈ Judge 71.1%** → niezależny, nie-LLM-owy sygnał potwierdza sędziego",
    "**Przykład (C1):** gold *Elizabeth Rosemond Taylor*, model *„Elizabeth Taylor was a British-American actress…\"* — merytorycznie **dobrze**, EM = 0",
    "→ EM/F1 mierzą formę, nie wiedzę. Stąd konieczność sędziego semantycznego.",
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.0), size=17)

prs.save(OUT)
print(f"Zapisano {OUT}  ({len(prs.slides)} slajdów)")
